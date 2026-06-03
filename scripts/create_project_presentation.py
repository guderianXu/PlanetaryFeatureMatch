#!/usr/bin/env python3
"""Create a project-introduction PPTX for PlanetaryFeatureMatch."""

from __future__ import annotations

import csv
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = PROJECT_ROOT / "runs" / "project_presentation_20260526"
ASSET_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "PlanetaryFeatureMatch_project_intro_20260526.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

NAVY = RGBColor(22, 43, 58)
TEAL = RGBColor(0, 132, 132)
BLUE = RGBColor(47, 92, 160)
GREEN = RGBColor(72, 143, 99)
ORANGE = RGBColor(210, 126, 54)
RED = RGBColor(184, 69, 67)
GRAY = RGBColor(103, 112, 120)
LIGHT_BG = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def set_run_font(run, size: float, bold: bool = False, color: RGBColor = NAVY) -> None:
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.62))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.name = FONT_CN
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.86), Inches(11.7), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.runs[0].font.name = FONT_CN
        sp.runs[0].font.size = Pt(12)
        sp.runs[0].font.color.rgb = GRAY
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.17), Inches(12.25), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, text: str = "PlanetaryFeatureMatch | 2026-05-26") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.14), Inches(12.25), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.runs[0].font.name = FONT_EN
    p.runs[0].font.size = Pt(8.5)
    p.runs[0].font.color.rgb = GRAY


def add_bullets(slide, x: float, y: float, w: float, h: float, bullets: Iterable[str], size: float = 15) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.color.rgb = NAVY


def add_metric_card(slide, x: float, y: float, w: float, h: float, value: str, label: str, color: RGBColor = TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(215, 222, 228)
    shape.line.width = Pt(1)
    v = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.18), Inches(w - 0.32), Inches(0.42))
    vp = v.text_frame.paragraphs[0]
    vp.text = value
    vp.alignment = PP_ALIGN.CENTER
    vp.runs[0].font.name = FONT_EN
    vp.runs[0].font.size = Pt(22)
    vp.runs[0].font.bold = True
    vp.runs[0].font.color.rgb = color
    l = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.68), Inches(w - 0.28), Inches(h - 0.76))
    lp = l.text_frame.paragraphs[0]
    lp.text = label
    lp.alignment = PP_ALIGN.CENTER
    lp.runs[0].font.name = FONT_CN
    lp.runs[0].font.size = Pt(10)
    lp.runs[0].font.color.rgb = GRAY


def add_box(slide, x: float, y: float, w: float, h: float, text: str, color: RGBColor, size: float = 12) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.name = FONT_CN
    p.runs[0].font.size = Pt(size)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = GRAY) -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass


def add_picture_fit(slide, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    if not image_path.exists():
        return
    with Image.open(image_path) as img:
        iw, ih = img.size
    target_w = Inches(w)
    target_h = Inches(h)
    scale = min(target_w / iw, target_h / ih)
    final_w = int(iw * scale)
    final_h = int(ih * scale)
    left = Inches(x) + int((target_w - final_w) / 2)
    top = Inches(y) + int((target_h - final_h) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=final_w, height=final_h)


def save_bar_chart(path: Path, labels: list[str], series: list[tuple[str, list[float], str]], ylabel: str) -> None:
    plt.figure(figsize=(10.5, 4.2), dpi=180)
    ax = plt.gca()
    count = len(labels)
    width = 0.8 / len(series)
    x_positions = list(range(count))
    for idx, (name, values, color) in enumerate(series):
        offsets = [x - 0.4 + width / 2 + idx * width for x in x_positions]
        ax.bar(offsets, values, width=width, label=name, color=color)
    ax.set_xticks(x_positions)
    ax.set_xticklabels(labels, rotation=18, ha="right", fontsize=9)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel(ylabel)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, loc="upper left", ncol=len(series))
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    plt.tight_layout()
    plt.savefig(path, transparent=False, facecolor="white")
    plt.close()


def save_line_chart(path: Path, labels: list[str], values: list[float], ylabel: str) -> None:
    plt.figure(figsize=(8.8, 3.2), dpi=180)
    ax = plt.gca()
    ax.plot(labels, values, marker="o", color="#008484", linewidth=2.5)
    ax.set_ylim(0, max(values) * 1.15 if values else 1.0)
    ax.set_ylabel(ylabel)
    ax.grid(axis="y", alpha=0.25)
    for x, y in zip(labels, values):
        ax.text(x, y, f"{y:.3f}", ha="center", va="bottom", fontsize=8)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    plt.tight_layout()
    plt.savefig(path, transparent=False, facecolor="white")
    plt.close()


def read_training_eval(path: Path) -> tuple[list[str], list[float], list[float]]:
    if not path.exists():
        return [], [], []
    with path.open(newline="", encoding="utf-8") as handle:
        rows = list(csv.DictReader(handle))
    if not rows:
        return [], [], []
    labels = []
    top1 = []
    top5 = []
    for row in rows:
        labels.append(row.get("phase", row.get("step", str(len(labels)))))
        if "top1" in row:
            top1.append(float(row["top1"]))
        elif "descriptor_top1" in row:
            top1.append(float(row["descriptor_top1"]))
        if "top5" in row:
            top5.append(float(row["top5"]))
        elif "descriptor_top5" in row:
            top5.append(float(row["descriptor_top5"]))
    return labels, top1, top5


def build_assets() -> dict[str, Path]:
    pure_labels = [
        "num-rot",
        "num-view",
        "num-comp",
        "ts-rot",
        "ts-view",
        "ts-comp",
    ]
    pure_precision = [0.933622, 0.370629, 0.223214, 0.902439, 0.192308, 0.138889]
    hybrid_precision = [0.997457, 0.969603, 0.930609, 0.996849, 0.985647, 0.982751]
    save_bar_chart(
        ASSET_DIR / "pure_vs_hybrid_precision.png",
        pure_labels,
        [
            ("Pure PFM", pure_precision, "#2F5CA0"),
            ("PFM + RootSIFT fallback", hybrid_precision, "#488F63"),
        ],
        "Precision",
    )

    cross_labels = ["num-view", "num-comp", "ts-view", "ts-comp"]
    pfm_sample = [0.361702, 0.122807, 0.112903, 0.104478]
    rootsift_h = [0.981952, 0.972522, 0.998124, 0.995624]
    save_bar_chart(
        ASSET_DIR / "traditional_cross_view_precision.png",
        cross_labels,
        [
            ("PFM guarded", pfm_sample, "#B84543"),
            ("RootSIFT + H-RANSAC", rootsift_h, "#008484"),
        ],
        "Precision",
    )

    rotated_labels = [
        "PFM raw",
        "PFM H",
        "AKAZE",
        "ORB",
        "RootSIFT r0.80",
        "RootSIFT r0.90",
        "LightGlue-SIFT",
    ]
    rotated_precision = [0.038408, 0.125, 0.993409, 0.995932, 1.0, 0.999471, 1.0]
    save_bar_chart(
        ASSET_DIR / "rotated_cross_view_precision.png",
        rotated_labels,
        [("rotated cross-view precision", rotated_precision, "#2F5CA0")],
        "Precision",
    )
    deep_rotation_labels = ["PFM", "LG-SIFT", "LG-ALIKED", "LG-DISK", "LG-SP", "LoFTR"]
    deep_rotation_precision = [1.0, 0.999349, 0.590909, 0.032520, 0.014493, 0.004886]
    save_bar_chart(
        ASSET_DIR / "deep_rotation_smoke_precision.png",
        deep_rotation_labels,
        [("same-image rotation smoke", deep_rotation_precision, "#008484")],
        "Precision",
    )
    p1_labels = ["num-view\nfixed64", "num-view\nfull-val", "ts-view\nfixed64", "ts-view\nfull-val"]
    p1_baseline = [0.370629, 0.443780, 0.192308, 0.200514]
    p1_probe = [0.334129, 0.351742, 0.117021, 0.208543]
    save_bar_chart(
        ASSET_DIR / "p1_descriptor_probe_precision.png",
        p1_labels,
        [
            ("current route", p1_baseline, "#2F5CA0"),
            ("P1 descriptor probe", p1_probe, "#B84543"),
        ],
        "Precision",
    )

    labels, top1, _ = read_training_eval(
        PROJECT_ROOT / "runs" / "cross_view_1024_rootsift_pseudo_r080t2_only_lr1e6_weakgroups_80_seed1234" / "training" / "eval_summary.csv"
    )
    if labels and top1:
        save_line_chart(ASSET_DIR / "descriptor_top1_example.png", labels, top1, "Top-1 retrieval")
    else:
        save_line_chart(ASSET_DIR / "descriptor_top1_example.png", ["before", "after"], [0.1004, 0.1027], "Top-1 retrieval")

    return {
        "pure_vs_hybrid": ASSET_DIR / "pure_vs_hybrid_precision.png",
        "traditional": ASSET_DIR / "traditional_cross_view_precision.png",
        "rotated": ASSET_DIR / "rotated_cross_view_precision.png",
        "deep_rotation": ASSET_DIR / "deep_rotation_smoke_precision.png",
        "p1_probe": ASSET_DIR / "p1_descriptor_probe_precision.png",
        "top1": ASSET_DIR / "descriptor_top1_example.png",
        "pure_num_comp": PROJECT_ROOT
        / "runs"
        / "cross_view_1024_keypointonly_multistate_lowcontrast_targetcontrast_postselect_0step_seed1234"
        / "visualizations"
        / "numeric"
        / "compound"
        / "01_source_000201_72_pair_002049.png",
        "pure_ts_view": PROJECT_ROOT
        / "runs"
        / "cross_view_1024_keypointonly_multistate_lowcontrast_targetcontrast_postselect_0step_seed1234"
        / "visualizations"
        / "timestamp"
        / "viewpoint"
        / "01_source_000123_20260514T144405909_NAS_PAN_L2b_pair_003819.png",
        "hybrid_ts_comp": PROJECT_ROOT
        / "runs"
        / "cross_view_1024_targetcontrast_rootsift_allgatezero_fallback_route_20260526"
        / "visualizations"
        / "timestamp"
        / "compound"
        / "01_source_000123_20260514T144405909_NAS_PAN_L2b_pair_003819.png",
    }


def setup_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(250, 251, 252)
    return slide


def create_deck() -> None:
    ensure_dirs()
    assets = build_assets()
    prs = setup_presentation()

    slide = blank_slide(prs)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.8), Inches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "PlanetaryFeatureMatch"
    p.runs[0].font.name = FONT_EN
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.92), Inches(11.5), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "行星影像深度特征提取与深度匹配算法项目介绍"
    sp.runs[0].font.name = FONT_CN
    sp.runs[0].font.size = Pt(26)
    sp.runs[0].font.color.rgb = TEAL
    add_bullets(
        slide,
        0.95,
        3.1,
        7.2,
        2.2,
        [
            "目标：在跨高度、跨视角、弱纹理行星影像中建立可训练、可评估的局部特征匹配流程",
            "当前状态：已形成 1024 cache 训练/验证/测试闭环；P1 retention probe 已完成但不进入 pure-PFM route",
            "汇报范围：模型结构、训练策略、性能结果、与传统/深度 matcher 对比、后续路线",
        ],
        size=15,
    )
    add_metric_card(slide, 9.0, 3.02, 2.8, 0.95, "6", "两类影像风格 x 三类 gate")
    add_metric_card(slide, 9.0, 4.1, 2.8, 0.95, "1024", "当前主数据集尺度")
    add_metric_card(slide, 9.0, 5.18, 2.8, 0.95, "pure / hybrid", "严格分离模型与外部 fallback")
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "任务背景：行星影像匹配难点")
    add_bullets(
        slide,
        0.75,
        1.45,
        5.7,
        4.8,
        [
            "输入影像可能来自不同拍摄高度，例如 1-155.tif 与 20260514T064636672_NAS_PAN_L2b.tif 属于不同成像尺度/高度。",
            "主要扰动：旋转、视角/高度变化、复合几何变化、光照差异、弱纹理和目标视图低对比度。",
            "评估拆为两类风格：numeric 与 timestamp/NAS；三类 gate：Rotate、Viewpoint、Compound。",
            "核心要求：不是只在固定样本上好看，而是按 train / val / test 划分训练，并用 full-val guard 防止过拟合。",
        ],
        size=15,
    )
    add_box(slide, 7.0, 1.55, 1.7, 0.65, "Rotate", BLUE)
    add_box(slide, 9.0, 1.55, 1.7, 0.65, "Viewpoint", TEAL)
    add_box(slide, 11.0, 1.55, 1.7, 0.65, "Compound", ORANGE)
    add_picture_fit(slide, assets["pure_ts_view"], 7.0, 2.45, 5.7, 3.6)
    add_footer(slide, "Source: README.md, current 1024 split/eval runs")

    slide = blank_slide(prs)
    add_title(slide, "端到端流程：从合成缓存到路由评估")
    boxes = [
        (0.75, 1.62, "原始行星影像\nnumeric / timestamp"),
        (3.0, 1.62, "Synthetic Pair Cache\nview_a, view_b, warp, mask"),
        (5.45, 1.62, "PFM 训练\nbackbone + heads + matcher"),
        (7.9, 1.62, "Validation Calibration\ncheckpoint / margin / score"),
        (10.35, 1.62, "Fixed Test + Full-Val Guard\npure 与 hybrid 分离"),
    ]
    colors = [NAVY, BLUE, TEAL, GREEN, ORANGE]
    for (x, y, text), color in zip(boxes, colors):
        add_box(slide, x, y, 1.85, 0.9, text, color, size=10)
    for x in [2.6, 5.05, 7.5, 9.95]:
        add_arrow(slide, x, 2.05, x + 0.35, 2.05, GRAY)
    add_bullets(
        slide,
        1.0,
        3.2,
        11.1,
        2.5,
        [
            "训练侧：自监督 warp CE、graph matching loss、RootSIFT pseudo-label、heatmap-only 蒸馏、target-quality gated specialist。",
            "评估侧：每组单独 calibration，使用 validation 选择参数；fixed test 输出随机两对可视化；关键新路线必须追加 full-val guard。",
            "策略侧：PFM learned route、传统 matcher teacher、外部 fallback hybrid 指标分开展示，避免把不同层级收益混成一个结论。",
        ],
        size=15,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "深度特征提取模型：Backbone + Sparse/Dense Heads")
    add_box(slide, 0.85, 1.55, 2.1, 0.72, "Backbone\n4-stage CNN", NAVY)
    add_box(slide, 3.55, 1.05, 2.55, 0.78, "SparseHead\nheatmap / descriptor / scale / orientation / affine", TEAL, size=10)
    add_box(slide, 3.55, 2.08, 2.55, 0.78, "DenseHead\nlocal correlation + offsets", BLUE, size=10)
    add_box(slide, 6.85, 1.55, 2.1, 0.72, "FeatureSet\nsparse + semi-dense", GREEN)
    add_box(slide, 9.75, 1.55, 2.1, 0.72, "Matcher / Eval\ngeometry + metrics", ORANGE)
    add_arrow(slide, 2.95, 1.91, 3.45, 1.43)
    add_arrow(slide, 2.95, 1.91, 3.45, 2.45)
    add_arrow(slide, 6.15, 1.43, 6.78, 1.88)
    add_arrow(slide, 6.15, 2.45, 6.78, 1.88)
    add_arrow(slide, 8.98, 1.91, 9.65, 1.91)
    add_bullets(
        slide,
        0.95,
        3.22,
        5.7,
        3.0,
        [
            "Backbone 默认容量：base_channels=32，四级下采样卷积 + refinement。",
            "SparseHead 共享 context tower，输出关键点 heatmap、128 维描述子、尺度、方向、仿射形状。",
            "descriptor 分支加入多尺度上下文、anisotropic viewpoint context、orientation alignment、dilated context。",
        ],
        size=14,
    )
    add_bullets(
        slide,
        7.0,
        3.22,
        5.4,
        3.0,
        [
            "DenseHead 融合两图 feature、差分、坐标通道和半径 4 的 local correlation，输出半稠密 confidence 与局部 offset。",
            "训练损失覆盖 repeatability、descriptor CE、graph matching CE、offset Smooth L1、confidence BCE。",
            "当前 1024 实验主线以稀疏匹配/learned keypoint score 为核心，dense 输出保留为后续半稠密扩展。",
        ],
        size=14,
    )
    add_footer(slide, "Source: README.md, python/pfm_model.py")

    slide = blank_slide(prs)
    add_title(slide, "特征匹配深度学习模型：Graph Matcher")
    add_box(slide, 0.9, 1.5, 2.0, 0.7, "A keypoints\nD=128 descriptors", BLUE)
    add_box(slide, 0.9, 2.6, 2.0, 0.7, "B keypoints\nD=128 descriptors", BLUE)
    add_box(slide, 3.55, 1.45, 2.2, 0.82, "Descriptor Projection\n+ radial keypoint embedding", TEAL, size=10)
    add_box(slide, 6.4, 1.45, 2.35, 0.82, "Self Attention\nCross Attention\nFFN x 4", NAVY, size=10)
    add_box(slide, 9.45, 1.45, 2.3, 0.82, "Pair logits\n+ dustbin row/col", GREEN, size=10)
    add_arrow(slide, 2.95, 1.85, 3.45, 1.85)
    add_arrow(slide, 5.82, 1.85, 6.3, 1.85)
    add_arrow(slide, 8.82, 1.85, 9.35, 1.85)
    add_bullets(
        slide,
        0.9,
        4.0,
        5.7,
        2.5,
        [
            "关键点坐标不再直接使用 x/y 绝对位置，而是半径与半径平方 embedding，减少同屏位置捷径。",
            "每层同时做图内 self-attention 和跨图 cross-attention，使描述子匹配具备上下文推理能力。",
            "输出包含 dustbin，用于显式表示不可匹配点；推理再做 mutual nearest 过滤。",
        ],
        size=14,
    )
    add_bullets(
        slide,
        7.0,
        4.0,
        5.5,
        2.5,
        [
            "图匹配 loss 用 decoded sparse keypoints、warp positives、deterministic negatives 和 dustbin 监督。",
            "实验显示 graph loss 可在训练中收敛，但跨高度 viewpoint/compound 的 sparse precision 仍是主要瓶颈。",
            "当前改进重心转向更安全的 keypoint/descriptor 蒸馏与路由约束，而不是盲目扩大 heatmap positives。",
        ],
        size=14,
    )
    add_footer(slide, "Source: python/pfm_model.py, findings.md")

    slide = blank_slide(prs)
    add_title(slide, "训练与蒸馏策略：有效信号与负结果")
    add_bullets(
        slide,
        0.85,
        1.45,
        5.95,
        5.2,
        [
            "自监督 synthetic warp：提供基础 descriptor/graph matching 监督，但对真实跨高度 viewpoint 的泛化有限。",
            "RootSIFT/H-RANSAC pseudo-label：descriptor-only 小步训练可出现轻微 retrieval 信号，但 sparse matching 容易过激活。",
            "Keypoint heatmap-only 蒸馏：不破坏 descriptor validation，已产生 numeric/compound 与 timestamp/compound 的局部提升。",
            "Target-view local contrast gate：在 timestamp/compound 上提高 precision，并通过 lowcontrast specialist 保住 fixed-test 5 个 correct。",
        ],
        size=14,
    )
    add_picture_fit(slide, assets["top1"], 7.1, 1.55, 5.45, 2.25)
    add_metric_card(slide, 7.25, 4.35, 2.3, 0.86, "5/36", "pure PFM timestamp/compound fixed-test")
    add_metric_card(slide, 9.95, 4.35, 2.3, 0.86, "12/126", "same route full-val guard")
    add_metric_card(slide, 7.25, 5.35, 2.3, 0.86, "No", "broad gate-zero heatmap rejected")
    add_metric_card(slide, 9.95, 5.35, 2.3, 0.86, "P1 reject", "descriptor-retention probe not routeable", RED)
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "当前 pure-PFM 固定测试性能")
    add_picture_fit(slide, assets["pure_vs_hybrid"], 0.75, 1.35, 7.1, 3.2)
    add_bullets(
        slide,
        8.05,
        1.45,
        4.65,
        3.35,
        [
            "当前最佳 pure-PFM route：lowcontrast + target-contrast postselected route。",
            "Rotate 两组已较稳定：numeric 647/693，timestamp 555/615。",
            "Viewpoint 与 Compound 仍弱，尤其 timestamp/viewpoint 20/104、timestamp/compound 5/36。",
            "图中 hybrid 仅作外部 fallback 对照，不能算 learned PFM 进步。",
        ],
        size=14,
    )
    add_picture_fit(slide, assets["pure_num_comp"], 0.85, 4.78, 5.7, 1.85)
    add_picture_fit(slide, assets["pure_ts_view"], 6.95, 4.78, 5.7, 1.85)
    add_footer(slide, "Source: runs/cross_view_1024_keypointonly_multistate_lowcontrast_targetcontrast_postselect_0step_seed1234")

    slide = blank_slide(prs)
    add_title(slide, "与传统匹配算法对比：RootSIFT/H-RANSAC 仍是强 teacher")
    add_picture_fit(slide, assets["traditional"], 0.75, 1.35, 7.0, 3.35)
    add_bullets(
        slide,
        8.0,
        1.42,
        4.65,
        4.4,
        [
            "真实 cross-view 样本上，PFM guarded summary 在 viewpoint/compound 弱组只有约 0.10-0.36 precision。",
            "RootSIFT-FLANN-ratio + Homography RANSAC 在四个弱组达到 0.9725-0.9981 precision。",
            "结论：信号并非几何上不可得，传统局部特征已经能在大量样本上找到高置信对应。",
            "这也是后续 pseudo-label / teacher mining 的主要依据。",
        ],
        size=14,
    )
    add_footer(slide, "Source: runs/cross_view_traditional_matcher_comparison_agent3/summary.md")

    slide = blank_slide(prs)
    add_title(slide, "与其他深度匹配算法对比：LightGlue-SIFT 很强，LoFTR/learned local feature 不稳定")
    add_picture_fit(slide, assets["rotated"], 0.65, 1.28, 6.7, 2.45)
    add_picture_fit(slide, assets["deep_rotation"], 0.65, 4.02, 6.7, 2.45)
    add_bullets(
        slide,
        7.85,
        1.42,
        4.75,
        4.9,
        [
            "旋转 cross-view 对比中，LightGlue-SIFT + Homography 达到 2089/2089，precision=1.0。",
            "同图 90/180/270 smoke 中，LightGlue-SIFT 接近满分；LightGlue-ALIKED 为 221/374，DISK、SuperPoint 和 LoFTR 明显不稳。",
            "RootSIFT r0.90/H2 为 1890/1891，precision=0.999471；SIFT/AKAZE/ORB 经几何过滤也很强。",
            "PFM raw 只有 83/2161，PFM-Homography 1/8，说明当前 learned matcher 置信度和 keypoint 选择仍不可靠。",
            "SuperGlue 本地缺少 match_pairs 入口，记录为 unavailable，没有纳入数值横向排名。",
        ],
        size=14,
    )
    add_footer(slide, "Source: matcher_algorithm_iteration_agent11, rotation_matcher_comparison_agent")

    slide = blank_slide(prs)
    add_title(slide, "Hybrid/Fallback 路线：部署价值高，但必须和 pure-PFM 分开")
    add_picture_fit(slide, assets["hybrid_ts_comp"], 0.75, 1.35, 6.3, 3.0)
    add_bullets(
        slide,
        7.4,
        1.45,
        5.0,
        3.3,
        [
            "All-gate-zero RootSIFT fallback 在 fixed-test 中整体达到 61437/62255，precision=0.986860。",
            "Stage10 full-val 证明更严格 r0.75/H2 fallback 精度更高：183662/184456，precision=0.995695。",
            "Stage11 fixed-test r0.75/H2 将整体 precision 提高到 0.988985，wrong 从 818 降到 650，但 correct 支撑下降。",
            "这些都是外部 matcher fallback，不是 PFM 模型本身学到的能力。",
        ],
        size=14,
    )
    add_metric_card(slide, 0.95, 5.05, 2.4, 0.9, "0.9869", "r0.80/H2 fixed-test hybrid")
    add_metric_card(slide, 3.75, 5.05, 2.4, 0.9, "0.988985", "r0.75/H2 fixed-test hybrid")
    add_metric_card(slide, 6.55, 5.05, 2.4, 0.9, "0.9957", "r0.75/H2 full-val fallback")
    add_footer(slide, "Source: Stage10/Stage11 matcher sidecars")

    slide = blank_slide(prs)
    add_title(slide, "负结果复盘：为什么 broad gate-zero heatmap 不能继续扩")
    add_bullets(
        slide,
        0.9,
        1.45,
        6.0,
        4.8,
        [
            "外部 fallback 作为推理候选很准，但并不等价于“所有 fallback 点都是安全 heatmap positive”。",
            "Gate-zero pair 本来是 PFM 选择 abstain 的区域；直接加入大量 positive 会驱动 heatmap 在弱纹理/低置信区域过激活。",
            "负结果特征：fixed-test compound 有小幅提升，但 full-val precision 崩塌，说明不是泛化提升。",
            "结论：后续不能再做 dense positive-only heatmap；需要小 cap、source balance、hard negatives 与 full-val guard。",
        ],
        size=15,
    )
    add_metric_card(slide, 7.45, 1.65, 2.35, 0.95, "0.212 -> 0.0846", "numeric/compound full-val")
    add_metric_card(slide, 10.0, 1.65, 2.35, 0.95, "0.0952 -> 0.0072", "timestamp/compound full-val")
    add_metric_card(slide, 7.45, 3.05, 2.35, 0.95, "6844", "broad gate-zero labels")
    add_metric_card(slide, 10.0, 3.05, 2.35, 0.95, "115", "kept train pairs")
    add_metric_card(slide, 7.45, 4.45, 4.9, 0.95, "禁止 F1 / F2", "不再 broad all-gate-zero heatmap，不用 r0.80/H2 做最大支撑 teacher", RED)
    add_footer(slide, "Source: gatezero checkpoint eval, Stage12 policy diagnostic")

    slide = blank_slide(prs)
    add_title(slide, "最新 P1 结果与下一步：必须保留 abstention")
    add_picture_fit(slide, assets["p1_probe"], 0.75, 1.32, 6.55, 2.7)
    add_bullets(
        slide,
        0.9,
        4.25,
        6.25,
        2.2,
        [
            "P1 使用 train-only RootSIFT r0.75/H2 + warp truth，62 pairs / 744 labels，wrong=0。",
            "both-viewpoint P1 让 numeric/viewpoint full-val precision 从 0.4438 降到 0.3517，属于 activation growth。",
            "timestamp-only P1 只在 full-val 有弱正信号 0.2005 -> 0.2085，但 fixed64 从 0.1923 降到 0.1170。",
        ],
        size=13.5,
    )
    add_bullets(
        slide,
        7.4,
        1.45,
        4.9,
        5.1,
        [
            "决策：P1 descriptor checkpoint 不加入 pure-PFM routing。",
            "后续训练不能只加 positive correspondence；需要显式 negative/abstention/retention 约束，避免把原本应拒绝的 pair 激活。",
            "下一步候选：warp-aware hard negatives、score gate 保守校准、或从当前 route 的 zero/low-support pair 构造非激活监督。",
            "三条指标线继续分离：pure-PFM、teacher/pseudo-label、hybrid fallback。",
            "任何新 checkpoint 先通过 fixed64 与 full-val non-regression，再讨论 compound microprobe。",
        ],
        size=14,
    )
    add_footer(slide, "Source: P1 descriptor probe summaries, Stage12 policy diagnostic")

    slide = blank_slide(prs)
    add_title(slide, "关键产物与复现实验路径")
    add_bullets(
        slide,
        0.85,
        1.35,
        11.8,
        5.7,
        [
            "当前最佳 pure-PFM route：runs/cross_view_1024_keypointonly_multistate_lowcontrast_targetcontrast_postselect_0step_seed1234/",
            "传统 matcher cross-view 对比：runs/cross_view_traditional_matcher_comparison_agent3/",
            "旋转 cross-view 深度/传统 matcher 对比：runs/matcher_algorithm_iteration_agent11/",
            "高精度 hybrid full-val 验证：runs/matcher_algorithm_iteration_agent14_stage10/",
            "固定测试 r0.75/H2 hybrid：runs/matcher_algorithm_iteration_agent14_stage11/",
            "训练安全策略诊断：runs/matcher_algorithm_iteration_agent14_stage12/",
            "P1 viewpoint descriptor probe：runs/cross_view_1024_p1_viewpoint_retention_desc_lr5e7_b2_80_seed1234/",
            "P1 timestamp/viewpoint-only probe：runs/cross_view_1024_p1_timestamp_viewpoint_retention_desc_lr5e7_b2_60_seed1234/",
            "timestamp/viewpoint quality gate 诊断：runs/timestamp_viewpoint_quality_gate_diagnostic_20260526/",
        ],
        size=12.8,
    )
    add_footer(slide)

    prs.save(PPTX_PATH)


def read_fixed_comparison_summary() -> dict[tuple[str, str, str], dict[str, str]]:
    path = PROJECT_ROOT / "对比文档" / "summary.csv"
    if not path.exists():
        return {}
    with path.open(newline="", encoding="utf-8") as handle:
        return {
            (row["style"], row["gate"], row["algorithm"]): row
            for row in csv.DictReader(handle)
        }


def build_simple_assets() -> dict[str, Path]:
    assets = build_assets()
    fixed = read_fixed_comparison_summary()
    labels = ["num-rot", "num-view", "num-comp", "ts-rot", "ts-view", "ts-comp"]
    groups = [
        ("numeric", "rotate"),
        ("numeric", "viewpoint"),
        ("numeric", "compound"),
        ("timestamp", "rotate"),
        ("timestamp", "viewpoint"),
        ("timestamp", "compound"),
    ]

    def precisions(algorithm: str) -> list[float]:
        values = []
        for style, gate in groups:
            row = fixed.get((style, gate, algorithm))
            values.append(float(row["precision"]) if row else 0.0)
        return values

    if fixed:
        save_bar_chart(
            ASSET_DIR / "fixed12_match_precision.png",
            labels,
            [
                ("PFM-current", precisions("PlanetaryFeatureMatch-current"), "#B84543"),
                ("RootSIFT-r0.90/H2", precisions("RootSIFT-r0.90-Ht2"), "#008484"),
                ("LightGlue-SIFT/H3", precisions("LightGlue-SIFT-Ht3"), "#488F63"),
            ],
            "Precision",
        )
    assets["fixed12_precision"] = ASSET_DIR / "fixed12_match_precision.png"
    compare_root = PROJECT_ROOT / "对比文档" / "figures"
    assets["cmp_pfm_num_view"] = compare_root / "pfm" / "numeric" / "viewpoint" / "01" / "01_source_000201_72_pair_002049.png"
    assets["cmp_lg_num_view"] = compare_root / "other_models" / "numeric" / "viewpoint" / "01" / "LightGlue-SIFT-Ht3.png"
    assets["cmp_pfm_ts_comp"] = compare_root / "pfm" / "timestamp" / "compound" / "01" / "01_source_000123_20260514T144405909_NAS_PAN_L2b_pair_003819.png"
    assets["cmp_root_ts_comp"] = compare_root / "other_models" / "timestamp" / "compound" / "01" / "RootSIFT-r0.90-Ht2.png"
    return assets


def create_deck() -> None:
    """Create a concise deck focused on this project's model and fixed-pair match results."""
    ensure_dirs()
    assets = build_simple_assets()
    fixed = read_fixed_comparison_summary()
    prs = setup_presentation()

    slide = blank_slide(prs)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.88), Inches(11.8), Inches(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = "PlanetaryFeatureMatch"
    p.runs[0].font.name = FONT_EN
    p.runs[0].font.size = Pt(42)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "行星影像深度特征提取与特征匹配模型简介"
    sp.runs[0].font.name = FONT_CN
    sp.runs[0].font.size = Pt(25)
    sp.runs[0].font.color.rgb = TEAL
    add_bullets(
        slide,
        0.95,
        3.0,
        7.0,
        2.4,
        [
            "项目目标：面向不同拍摄高度、不同视角和弱纹理行星影像，学习稳定局部特征与匹配关系。",
            "当前版本：已完成 1024 synthetic cache 的 train / val / test 流程，按两种影像风格与三类 gate 评估。",
            "本汇报重点：先介绍当前模型，再展示固定 6 组 x 2 个 pair 上与其他方法的匹配效果对比。",
        ],
        size=15,
    )
    add_metric_card(slide, 9.0, 3.05, 2.8, 0.95, "1024", "主实验缓存尺寸")
    add_metric_card(slide, 9.0, 4.12, 2.8, 0.95, "6 x 2", "固定可视化对比样本")
    add_metric_card(slide, 9.0, 5.2, 2.8, 0.95, "pure PFM", "与外部 fallback 分开报告")
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "数据与评估口径")
    add_bullets(
        slide,
        0.85,
        1.45,
        6.1,
        4.8,
        [
            "两类影像风格：numeric 文件名影像、timestamp/NAS 影像；它们代表不同命名来源和成像风格。",
            "三类 gate：Rotate、Viewpoint、Compound，分别考察旋转、跨高度/视角、复合几何变化。",
            "之前生成的可视化样本固定为 6 组 x 每组 2 个 pair，共 12 个 pair；对比文档复用同一批 pair。",
            "匹配正确性按 warp 真值计算，阈值为 5 px；外部方法的图和 CSV 已放入 `对比文档/`。",
        ],
        size=15,
    )
    add_picture_fit(slide, assets["pure_ts_view"], 7.25, 1.55, 5.0, 3.25)
    add_footer(slide, "Samples: 对比文档/fixed_pairs.csv")

    slide = blank_slide(prs)
    add_title(slide, "当前项目模型：深度特征提取")
    add_box(slide, 0.85, 1.55, 2.1, 0.72, "Input image\n1-channel", NAVY)
    add_box(slide, 3.45, 1.55, 2.15, 0.72, "Backbone\nmulti-stage CNN", BLUE)
    add_box(slide, 6.1, 1.10, 2.55, 0.78, "SparseHead\nheatmap + descriptor", TEAL, size=10)
    add_box(slide, 6.1, 2.12, 2.55, 0.78, "DenseHead\nconfidence + offset", GREEN, size=10)
    add_box(slide, 9.5, 1.55, 2.35, 0.72, "Feature outputs\nsparse / semi-dense", ORANGE, size=10)
    add_arrow(slide, 2.98, 1.91, 3.35, 1.91)
    add_arrow(slide, 5.65, 1.91, 6.0, 1.48)
    add_arrow(slide, 5.65, 1.91, 6.0, 2.48)
    add_arrow(slide, 8.72, 1.48, 9.42, 1.88)
    add_arrow(slide, 8.72, 2.48, 9.42, 1.88)
    add_bullets(
        slide,
        0.95,
        3.55,
        11.4,
        2.3,
        [
            "Backbone 提取共享多尺度特征；SparseHead 输出关键点 heatmap、128 维描述子、尺度/方向/仿射形状。",
            "DenseHead 保留半稠密置信度和局部 offset 能力，当前 1024 迭代主要以稀疏关键点与描述子匹配为主。",
            "模型训练来自 synthetic cache 的 warp 监督，并逐步加入 RootSIFT/LightGlue 等外部高置信对应作为 teacher 信号。",
        ],
        size=15,
    )
    add_footer(slide, "Source: README.md, python/pfm_model.py")

    slide = blank_slide(prs)
    add_title(slide, "当前项目模型：深度匹配模块")
    add_box(slide, 0.95, 1.45, 2.1, 0.78, "A image\nkeypoints + descriptors", BLUE, size=10)
    add_box(slide, 0.95, 2.55, 2.1, 0.78, "B image\nkeypoints + descriptors", BLUE, size=10)
    add_box(slide, 3.75, 1.95, 2.3, 0.82, "Descriptor projection\n+ keypoint embedding", TEAL, size=10)
    add_box(slide, 6.65, 1.95, 2.3, 0.82, "Self/Cross Attention\nGraph reasoning", NAVY, size=10)
    add_box(slide, 9.55, 1.95, 2.2, 0.82, "Match logits\n+ dustbin", GREEN, size=10)
    add_arrow(slide, 3.1, 1.84, 3.65, 2.33)
    add_arrow(slide, 3.1, 2.94, 3.65, 2.33)
    add_arrow(slide, 6.12, 2.36, 6.55, 2.36)
    add_arrow(slide, 9.0, 2.36, 9.45, 2.36)
    add_bullets(
        slide,
        0.95,
        3.75,
        11.3,
        2.4,
        [
            "Graph Matcher 使用描述子和关键点位置编码做多层自注意力与交叉注意力，输出匹配 logits。",
            "dustbin 表示不可匹配点，推理阶段再结合 mutual nearest、margin 和局部几何过滤得到最终 sparse matches。",
            "当前问题不是流程不可运行，而是跨高度 viewpoint/compound 下 learned descriptor/keypoint 的 precision 与召回仍不稳定。",
        ],
        size=15,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "训练与评估闭环")
    boxes = [
        (0.8, "1024 synthetic cache\ntrain / val / test"),
        (3.25, "PFM training\nwarp + pseudo labels"),
        (5.7, "Validation calibration\ncheckpoint / margin / score"),
        (8.15, "Fixed test\n6 groups"),
        (10.55, "Visual report\n2 pairs per group"),
    ]
    colors = [NAVY, BLUE, TEAL, GREEN, ORANGE]
    for (x, text), color in zip(boxes, colors):
        add_box(slide, x, 1.65, 1.85, 0.95, text, color, size=10)
    for x in [2.72, 5.18, 7.62, 10.03]:
        add_arrow(slide, x, 2.12, x + 0.35, 2.12)
    add_bullets(
        slide,
        0.95,
        3.35,
        11.3,
        2.55,
        [
            "每次训练先在 train split 上更新，再用 validation 选择 route 参数，最后在 test split 上报告六组指标。",
            "runs 目录保留每组随机两个匹配对的图，方便直接观察成功/失败模式。",
            "外部算法对比使用同样 12 个 pair，结果集中在 `对比文档/`，避免 PPT 里混入算法原理介绍。",
        ],
        size=15,
    )
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "当前 pure-PFM：完整 fixed-test 六组结果")
    add_picture_fit(slide, assets["pure_vs_hybrid"], 0.75, 1.35, 7.1, 3.2)
    add_bullets(
        slide,
        8.05,
        1.45,
        4.65,
        3.6,
        [
            "当前最佳 pure-PFM route：lowcontrast + target-contrast postselected route。",
            "Rotate 两组已经较稳定：numeric 647/693，timestamp 555/615。",
            "Viewpoint 和 Compound 仍是主要瓶颈，尤其 timestamp/viewpoint 与 timestamp/compound。",
            "hybrid/fallback 只作为部署或 teacher 参考，不计入 pure-PFM 模型指标。",
        ],
        size=14,
    )
    add_picture_fit(slide, assets["pure_num_comp"], 0.85, 4.8, 5.7, 1.8)
    add_picture_fit(slide, assets["pure_ts_view"], 6.95, 4.8, 5.7, 1.8)
    add_footer(slide, "Source: current pure-PFM fixed-test route")

    slide = blank_slide(prs)
    add_title(slide, "固定 12 对：与其他方法的匹配结果对比")
    add_picture_fit(slide, assets["fixed12_precision"], 0.7, 1.35, 7.25, 3.35)
    add_bullets(
        slide,
        8.15,
        1.45,
        4.55,
        4.3,
        [
            "这里不是介绍其他模型结构，只比较同一批 12 个 pair 的匹配效果。",
            "外部方法包括 SIFT、RootSIFT、ORB、AKAZE，以及 LightGlue-SIFT；SuperGlue 本地依赖不可用。",
            "这批样例上，RootSIFT/LightGlue-SIFT 在多数 group 能输出大量正确匹配；PFM 在 viewpoint/compound 样例上明显不足。",
            "完整 CSV、每个算法的匹配图都在 `对比文档/`。",
        ],
        size=14,
    )
    add_footer(slide, "Source: 对比文档/summary.csv")

    slide = blank_slide(prs)
    add_title(slide, "匹配图示例：numeric/viewpoint")
    add_picture_fit(slide, assets["cmp_pfm_num_view"], 0.75, 1.45, 5.8, 4.55)
    add_picture_fit(slide, assets["cmp_lg_num_view"], 6.85, 1.45, 5.8, 4.55)
    add_metric_card(slide, 1.35, 6.18, 4.6, 0.68, "PFM-current: 0/0", "当前模型在该样例没有有效匹配", RED)
    row = fixed.get(("numeric", "viewpoint", "LightGlue-SIFT-Ht3"), {})
    lg_text = f"LightGlue-SIFT: {row.get('correct', '500')}/{row.get('inlier_matches', '500')}"
    add_metric_card(slide, 7.45, 6.18, 4.6, 0.68, lg_text, "同一组两个样例聚合结果", GREEN)
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "匹配图示例：timestamp/compound")
    add_picture_fit(slide, assets["cmp_pfm_ts_comp"], 0.75, 1.45, 5.8, 4.55)
    add_picture_fit(slide, assets["cmp_root_ts_comp"], 6.85, 1.45, 5.8, 4.55)
    add_metric_card(slide, 1.35, 6.18, 4.6, 0.68, "PFM-current: 0/0", "当前模型在该样例没有有效匹配", RED)
    row = fixed.get(("timestamp", "compound", "RootSIFT-r0.90-Ht2"), {})
    rs_text = f"RootSIFT-r0.90/H2: {row.get('correct', '291')}/{row.get('inlier_matches', '291')}"
    add_metric_card(slide, 7.45, 6.18, 4.6, 0.68, rs_text, "同一组两个样例聚合结果", GREEN)
    add_footer(slide)

    slide = blank_slide(prs)
    add_title(slide, "结论与下一步")
    add_bullets(
        slide,
        0.95,
        1.4,
        11.6,
        5.25,
        [
            "当前模型结构和评估闭环已经可运行：能按 1024 cache 的 train/val/test 流程训练、校准、评估和画图。",
            "固定 12 对样例说明：外部方法能在同一图像对上找到大量正确匹配，因此数据中存在可学习信号。",
            "PFM 的主要差距在 learned keypoint/descriptor 对跨高度和复合视角的泛化，以及何时应该 abstain 的置信度控制。",
            "下一步训练重点：不要继续只加正样本 pseudo-label；需要加入 negative/abstention/retention 约束，让模型在增加正确匹配时不引入更多错误匹配。",
            "对比图和原始表格路径：`/home/xjw/code/deeplearning/PlanetaryFeatureMatch/对比文档/`。",
        ],
        size=16,
    )
    add_footer(slide)

    prs.save(PPTX_PATH)


def main() -> int:
    create_deck()
    print(PPTX_PATH)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
